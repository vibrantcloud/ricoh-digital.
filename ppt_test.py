from pathlib import Path
from pptx import Presentation
import copy
from time import sleep
from datetime import datetime
import re
import shutil



def convert_m_time(iterable):
    dt = datetime.fromtimestamp(iterable).strftime("%Y-%m-%d-%H:%M")
    return dt


def move_ppts():
    path = r"C:\Users\Umar.Hussain\Downloads"
    rep1 = [f for f in Path(path).glob("*Report 1*.*pptx")]
    rep2 = [f for f in Path(path).glob("*Report 2*.*pptx")]
    ppts = [rep1, rep2]
    return ppts


def get_max_unique_files():
    """currently two reports"""
    max_files = []
    files = move_ppts()
    max_files.append(max(files[0], key=lambda p: p.stat().st_ctime))
    max_files.append(max(files[1], key=lambda p: p.stat().st_ctime))
    return max_files


def move_files():
    folder_name = datetime.today().strftime("%y%m%d_%H%M%S_PPTS")
    trg_pptx_path = r"C:\Users\Umar.Hussain\OneDrive - Ricoh Europe PLC\Projects_FY19\MAG\Python Script\PPTX Converter"
    files_to_move = get_max_unique_files()
    for file in files_to_move:
        if Path(trg_pptx_path).joinpath(folder_name).is_dir() == False:
            Path.mkdir(Path(trg_pptx_path).joinpath(folder_name))
        else:
            pass
        shutil.copy(file, Path(trg_pptx_path).joinpath(folder_name))
        file.unlink()  # removes file in downloads.
    file_path = Path(trg_pptx_path).joinpath(folder_name)
    return file_path


def get_ppts_for_formatting():
    path = move_files()
    moved_ppts = [f for f in path.glob('*.pptx')]
    return moved_ppts,path


def format_power_points():    
    
    dimensions = {'height' : 7002000,
    'width' : 12193200,
    'top' : -57600,               
    'left' : 0 }
    
    ppts_to_format,path = get_ppts_for_formatting()
    
    
    for ppt in ppts_to_format:
        prs = Presentation(ppt)
    

        for slide in prs.slides:
            for shape in slide.shapes:
                shape.click_action.hyperlink.address = None
                shape.height = dimensions['height']
                shape.width = dimensions['width']
                shape.top = dimensions['top']
                shape.left = dimensions['left']
        delete_slide(prs,prs.slides[0])

        if Path(path).joinpath('Curated').is_dir() == False:
            Path.mkdir(Path(path).joinpath('Curated'))
        else:
            pass
        
        if Path(path).joinpath('Raw').is_dir() == False:
            Path.mkdir(Path(path).joinpath('Raw'))
        else:
            pass

        prs.save(Path(path).joinpath('Curated',f'{ppt.stem}_edited.pptx'))
        shutil.copy(ppt,Path(path).joinpath('Raw'))
        ppt.unlink()
    print("files saved.")
    

# Presentation Functions


def delete_slide(prs, slide):
    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def copy_slide(pres, pres1, index):
    source = pres.slides[index]
    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres1.slides.add_slide(blank_slide_layout)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")

    for key, value in source.part.rels.items():
        if not "notesSlide" in value.reltype:
            dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    return dest      

if __name__ == "__main__":
    format_power_points()